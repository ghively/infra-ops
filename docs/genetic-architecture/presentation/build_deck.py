#!/usr/bin/env python3
"""Build the infra-ops PoC presentation (16:9 PPTX).

Regenerate with:  python3 docs/genetic-architecture/presentation/build_deck.py
Output:           docs/genetic-architecture/presentation/infra-ops-poc.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
BG      = RGBColor(0x0E, 0x17, 0x26)   # deep navy
PANEL   = RGBColor(0x17, 0x24, 0x3B)   # card panel
PANEL2  = RGBColor(0x1E, 0x2E, 0x4A)   # lighter panel
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
TEXT    = RGBColor(0xF1, 0xF5, 0xF9)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
DARKTXT = RGBColor(0x0E, 0x17, 0x26)

FONT = "Calibri"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
_page = [0]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    _page[0] += 1
    return s


def txt(slide, x, y, w, h, lines, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4):
    """lines: str or list of (text, dict-overrides) or plain strings."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    first = True
    for ln in lines:
        ov = {}
        if isinstance(ln, tuple):
            ln, ov = ln
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        p.space_after = Pt(ov.get("space_after", space_after))
        r = p.add_run(); r.text = ln
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
    return box


def shape(slide, kind, x, y, w, h, fill=PANEL, line=None):
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp


def set_shape_text(sp, lines, size=12, color=TEXT, bold=False, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    if isinstance(lines, str):
        lines = [lines]
    first = True
    for ln in lines:
        ov = {}
        if isinstance(ln, tuple):
            ln, ov = ln
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space_after", 2))
        r = p.add_run(); r.text = ln
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)


def header(slide, kicker, title):
    txt(slide, Inches(0.6), Inches(0.32), Inches(12), Inches(0.35),
        kicker.upper(), size=12, color=TEAL, bold=True)
    txt(slide, Inches(0.6), Inches(0.62), Inches(12.2), Inches(0.8),
        title, size=28, bold=True)
    bar = shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.28), Inches(1.6), Pt(3), fill=AMBER)
    return bar


def footer(slide, note="infra-ops · governed multi-agent automation · PoC"):
    txt(slide, Inches(0.6), Inches(7.08), Inches(10), Inches(0.3), note, size=9, color=MUTED)
    txt(slide, Inches(12.3), Inches(7.08), Inches(0.6), Inches(0.3), str(_page[0]),
        size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent=TEAL, title_size=13, body_size=11):
    sp = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=PANEL)
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    tick = shape(slide, MSO_SHAPE.RECTANGLE, x, y + Inches(0.12), Pt(3.5), h - Inches(0.24), fill=accent)
    txt(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.3), Inches(0.4),
        title, size=title_size, bold=True, color=accent)
    if body:
        if isinstance(body, str):
            body = [body]
        txt(slide, x + Inches(0.18), y + Inches(0.46), w - Inches(0.32), h - Inches(0.55),
            body, size=body_size, color=TEXT, space_after=3)
    return sp


# ================================================================ 1 · TITLE
s = new_slide()
shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(5.95), W, Pt(2.5), fill=TEAL)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
    "PROOF OF CONCEPT · v0.14.0 · JUNE 2026", size=14, color=TEAL, bold=True)
txt(s, Inches(0.9), Inches(2.0), Inches(11.8), Inches(1.7),
    "infra-ops", size=60, bold=True)
txt(s, Inches(0.9), Inches(3.15), Inches(11.8), Inches(1.0),
    "Governed multi-agent automation for a PCI card-production estate",
    size=24, color=TEXT)
txt(s, Inches(0.9), Inches(4.1), Inches(11.8), Inches(1.2), [
    "Ansible · self-hosted GitLab CI/CD · Octopus Deploy · mixed Windows/Linux",
    "PCI DSS v4.0.1  +  PCI Card Production  +  PIN scope",
], size=15, color=MUTED, space_after=6)
txt(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.6),
    "Workflows · Architecture · Governance · Compliance · Harness strategy",
    size=14, color=AMBER, bold=True)

# ================================================================ 2 · PROBLEM / BET
s = new_slide()
header(s, "Why this exists", "The problem, and the bet we made")
card(s, Inches(0.6), Inches(1.65), Inches(6.0), Inches(4.9), "The problem", [
    "Infrastructure work in a card-production environment is slow by design — every",
    "change needs review evidence, compliance mapping, audit trail, and rollback plans.",
    "",
    "LLM agents are fast but unaccountable by default: free-form judgment,",
    "no audit trail, nothing stops a bad tool call.",
    "",
    "A QSA will never accept “the model was told to be careful.”",
], body_size=13)
card(s, Inches(6.85), Inches(1.65), Inches(5.9), Inches(4.9), "The bet", [
    "Make the agent admissible in a CDE by construction:",
    "",
    "•  Enforcement in code, not prompts — guardrails execute at the tool boundary",
    "•  Determinism where judgment is dangerous — merge decisions are computed",
    "•  Every action ledgered — append-only audit trail (PCI Req 10)",
    "•  Learning is human-gated — the system never silently changes its own behavior",
    "•  Hard zone separation — corporate (DSS) vs air-gapped HSA (CP + PIN)",
], accent=AMBER, body_size=13)
footer(s)

# ================================================================ 3 · WHAT'S BUILT
s = new_slide()
header(s, "Status", "What is built and wired today (v0.14.0)")
stats = [
    ("16", "specialist agents", "10 corporate + 6 air-gapped HSA, each with least-privilege tools & output contracts"),
    ("24", "skills", "lazy-loaded expertise — open SKILL.md standard, cross-vendor portable"),
    ("13", "code guardrails", "9 event-wired hooks + 4 CLI gates: DLP, zone routing, lint, audit, dual control"),
    ("8", "workflow commands", "scaffold, discover, drift-check, review, ingest, promote, rollback, preflight"),
    ("9", "audit collections", "append-only State Store + SIEM forwarder (PCI Req 10)"),
    ("14", "portable “genes”", "the architecture is captured as a framework-agnostic genome"),
]
for i, (n, label, desc) in enumerate(stats):
    x = Inches(0.6 + (i % 3) * 4.13)
    y = Inches(1.7 + (i // 3) * 2.5)
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.25), fill=PANEL)
    try: sp.adjustments[0] = 0.05
    except Exception: pass
    txt(s, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.9), n, size=44, bold=True, color=TEAL)
    txt(s, x + Inches(0.2), y + Inches(0.92), Inches(3.5), Inches(0.4), label, size=15, bold=True)
    txt(s, x + Inches(0.2), y + Inches(1.3), Inches(3.55), Inches(0.9), desc, size=11, color=MUTED)
txt(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.35),
    "Everything deterministic is a script with exit codes — runnable in CI with no model in the loop.",
    size=12, color=AMBER, bold=True)
footer(s)

# ================================================================ 4 · ARCHITECTURE
s = new_slide()
header(s, "Architecture", "Layered: a lean orchestrator over enforced boundaries")
layers = [
    ("Routing", "Orchestrator (lean) · 8 commands · 3 contexts — classify → delegate → assemble", TEAL),
    ("Specialists", "16 isolated agents, fresh context each call · machine-readable verdict contracts", TEAL),
    ("Quality gates", "3 parallel reviewers → deterministic merge gate (any BLOCK blocks · exit 0/1/3)", AMBER),
    ("Enforcement", "Tool-boundary hooks: PAN/DLP fail-closed · CHD sensitivity router · investigation gate · lint", RED),
    ("Knowledge & learning", "24 skills (lazy) · 13 path-scoped rules · governed instinct ledger (human-gated)", GREEN),
    ("Audit", "Append-only State Store (9 collections) · SIEM forwarding · PCI Req 10", MUTED),
]
y = Inches(1.62)
for name, desc, accent in layers:
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(8.5), Inches(0.74), fill=PANEL)
    try: sp.adjustments[0] = 0.12
    except Exception: pass
    shape(s, MSO_SHAPE.RECTANGLE, Inches(0.6), y + Inches(0.08), Pt(4), Inches(0.58), fill=accent)
    txt(s, Inches(0.78), y + Inches(0.05), Inches(2.0), Inches(0.6), name, size=13, bold=True, color=accent)
    txt(s, Inches(2.75), y + Inches(0.05), Inches(6.3), Inches(0.66), desc, size=11.5, color=TEXT,
        anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
# zones panel
zp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.35), Inches(1.62), Inches(3.4), Inches(4.84), fill=PANEL2)
try: zp.adjustments[0] = 0.05
except Exception: pass
txt(s, Inches(9.55), Inches(1.75), Inches(3.0), Inches(0.4), "TWO ZONES", size=12, bold=True, color=AMBER)
txt(s, Inches(9.55), Inches(2.15), Inches(3.05), Inches(1.9), [
    ("CORPORATE", {"bold": True, "color": TEAL, "size": 13}),
    "PCI DSS v4.0.1 · cloud inference allowed · propose-only (MRs, gated Dev deploy)",
], size=11, space_after=4)
txt(s, Inches(9.55), Inches(3.55), Inches(3.05), Inches(2.2), [
    ("HSA (AIR-GAPPED)", {"bold": True, "color": RED, "size": 13}),
    "PCI CP + PIN · local inference only · dual-control approvals · CPSA gates deployment",
    "",
    ("Crown jewels untouchable in both zones: PAN/CHD, keys, PINs, HSM config", {"color": MUTED, "size": 10}),
], size=11, space_after=4)
footer(s)

# ================================================================ 5 · CORE WORKFLOW
s = new_slide()
header(s, "Workflow showcase", "The authoring pipeline — LLM as witness, code as judge")
flow = [("Request", PANEL2, TEXT), ("Plan\n(opus)", PANEL, TEAL), ("Author\n(opus)", PANEL, TEAL),
        ("3 reviewers\nin parallel", PANEL, AMBER), ("Merge gate\n(script)", AMBER, DARKTXT)]
x = Inches(0.6)
for i, (label, fill, color) in enumerate(flow):
    kind = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
    sp = shape(s, kind, x, Inches(1.75), Inches(2.42), Inches(0.95), fill=fill)
    set_shape_text(sp, label.split("\n"), size=13, bold=True, color=color)
    x += Inches(2.47)
# reviewers detail
rv = [("playbook-reviewer", "correctness · idempotency · runs ansible-lint, syntax, --check --diff", "sonnet"),
      ("pci-compliance-reviewer", "PCI DSS v4.0.1 controls (Req 3/4/6/7/8/10/12)", "sonnet"),
      ("secrets-scanner", "static PAN / secret / key-material scan — never reproduces values", "haiku")]
for i, (n, d, m) in enumerate(rv):
    x = Inches(0.6 + i * 4.13)
    c = card(s, x, Inches(3.0), Inches(3.9), Inches(1.25), n, d, accent=AMBER, title_size=12, body_size=10.5)
    txt(s, x + Inches(2.9), Inches(3.08), Inches(0.9), Inches(0.3), m, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
txt(s, Inches(0.6), Inches(4.42), Inches(12.2), Inches(0.35),
    "Each reviewer returns a first-line token —  VERDICT: PASS | WARN | BLOCK  — the system's machine-readable ABI.",
    size=12.5, color=TEXT, bold=True)
outcomes = [
    ("ALL PASS / WARN → cleared (exit 0)", "Change merges. WARN is advisory and recorded.", GREEN),
    ("ANY BLOCK → revise (exit 1)", "Findings return to the author for ONE revision pass. Hard cap: 2 cycles.", AMBER),
    ("STILL BLOCKED → escalate (exit 3)", "Stops and hands open findings to a human. Never merges around a BLOCK. Missing verdict = BLOCK.", RED),
]
for i, (t, d, a) in enumerate(outcomes):
    card(s, Inches(0.6 + i * 4.13), Inches(4.95), Inches(3.9), Inches(1.7), t, d, accent=a, title_size=12, body_size=11)
footer(s)

# ================================================================ 6 · DAY-TO-DAY WORKFLOWS
s = new_slide()
header(s, "Workflow showcase", "Day-to-day operations — eight commands, all propose-only")
cmds = [
    ("/scaffold", "New role / module / pipeline from canonical templates; structure validated by script before any LLM touches it"),
    ("/preflight", "Fail-fast environment checklist: clean tree, no staged secrets, no placeholders"),
    ("/infra-discover", "Read-only estate discovery → knowledge/environment.md — the shared ground truth"),
    ("/drift-check", "ansible-playbook --check --diff against targets; reports drift, never auto-remediates"),
    ("/playbook-review", "Parallel correctness + PCI review of any MR or playbook with merged severity report"),
    ("/knowledge-ingest", "Classifies docs (PUBLIC → CHD-ADJACENT), indexes, answers scoping questions with citations"),
    ("/instinct-promote", "Human-gated learning: approver + confidence ≥ 0.7 + citations required"),
    ("/instinct-rollback", "First-class undo for learned behavior; HSA items need two approvers"),
]
for i, (n, d) in enumerate(cmds):
    x = Inches(0.6 + (i % 2) * 6.25)
    y = Inches(1.65 + (i // 2) * 1.25)
    card(s, x, y, Inches(6.0), Inches(1.12), n, d, accent=TEAL, title_size=13, body_size=10.5)
txt(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.35),
    "Hard rule #1 — propose, never dispose: the system edits code and opens MRs; promotion to test/staging/prod is always human-gated.",
    size=12, color=AMBER, bold=True)
footer(s)

# ================================================================ 7 · GOVERNANCE: ENFORCEMENT
s = new_slide()
header(s, "Governance", "Guardrails are code at the tool boundary — not instructions")
txt(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.4),
    "Every tool call the agent makes passes through these gates. Ignoring the prompt cannot bypass them.",
    size=13, color=MUTED)
hooks = [
    ("pan-egress-filter", "DLP: blocks Luhn-valid PAN / secrets in any tool input", "FAIL-CLOSED", RED),
    ("sensitivity-router", "Denies CHD-adjacent operations toward the cloud; routes to the local lane", "FAIL-CLOSED", RED),
    ("gateguard-fact-force", "Demands blast-radius + rollback facts before any mutating edit", "PRE-EDIT GATE", AMBER),
    ("yamllint + ansible-syntax", "Quality gates run automatically after every write", "POST-WRITE", AMBER),
    ("governance-capture / ledger", "Every tool execution appended to the audit store (async, non-blocking)", "ALWAYS ON", GREEN),
    ("dual-control gates", "HSA / compliance promotions and rollbacks require two named approvers", "TWO-PERSON", GREEN),
]
for i, (n, d, tag, a) in enumerate(hooks):
    y = Inches(2.05 + i * 0.78)
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.15), Inches(0.68), fill=PANEL)
    try: sp.adjustments[0] = 0.14
    except Exception: pass
    txt(s, Inches(0.8), y + Inches(0.05), Inches(3.2), Inches(0.55), n, size=13, bold=True, color=TEAL,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.05), y + Inches(0.05), Inches(6.6), Inches(0.58), d, size=11.5, anchor=MSO_ANCHOR.MIDDLE)
    tg = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.85), y + Inches(0.13), Inches(1.75), Inches(0.42), fill=a)
    set_shape_text(tg, tag, size=10, bold=True, color=DARKTXT)
footer(s)

# ================================================================ 8 · COMPLIANCE: AUDIT
s = new_slide()
header(s, "Compliance", "Audit trail by construction — PCI DSS Requirement 10")
card(s, Inches(0.6), Inches(1.65), Inches(6.0), Inches(3.3), "Unified State Store (append-only)", [
    "9 collections — every tool execution, decision, approval,",
    "promotion and rollback is a structured, schema-validated event:",
    "",
    "governanceEvents · observations · knowledgeBase · skillRuns",
    "skillVersions · sessions · workItems · decisions · installState",
    "",
    "Capture is asynchronous — auditing never slows the work.",
], body_size=12)
card(s, Inches(6.85), Inches(1.65), Inches(5.9), Inches(3.3), "SIEM forwarding", [
    "Governance events forward to the corporate SIEM endpoint",
    "(PCI CP §6.4 retention) — wiring is built, endpoint config is a",
    "deployment step on the roadmap.",
    "",
    "Verdicts, approvers, citations and confidence scores are part of",
    "the event payload: reviews are reconstructable end-to-end.",
], accent=GREEN, body_size=12)
card(s, Inches(0.6), Inches(5.1), Inches(12.15), Inches(1.55), "What an assessor sees", [
    "For any merged change: who proposed it (agent + session), the three reviewer verdicts with file:line findings, the computed gate decision",
    "with exit code, the revision cycles consumed, the human approver — and for learned behavior, the full promotion provenance with citations.",
], accent=AMBER, body_size=12.5)
footer(s)

# ================================================================ 9 · COMPLIANCE: LEARNING LOOP
s = new_slide()
header(s, "Compliance", "Self-improvement with a constitution — agents never self-promote")
steps = [("Observe", "hooks capture patterns from real work — automatic", TEAL),
         ("Curate", "knowledge-curator drafts instinct candidates with evidence + confidence", TEAL),
         ("Human gate", "approver required · confidence ≥ 0.7 · citations for compliance items", AMBER),
         ("Ledger", "zone-segmented YAML with full provenance (who, when, why, evidence)", GREEN),
         ("Rollback", "first-class undo · HSA & compliance items need TWO approvers", RED)]
x = Inches(0.6)
for i, (t, d, a) in enumerate(steps):
    kind = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
    sp = shape(s, kind, x, Inches(1.8), Inches(2.42), Inches(1.0), fill=PANEL if a != AMBER else AMBER)
    set_shape_text(sp, t, size=14, bold=True, color=DARKTXT if a == AMBER else a)
    txt(s, x + Inches(0.18), Inches(2.95), Inches(2.25), Inches(1.3), d, size=10.5, color=MUTED)
    x += Inches(2.47)
card(s, Inches(0.6), Inches(4.45), Inches(6.0), Inches(2.1), "Why it matters", [
    "“Self-improving agents” are the #1 emerging audit risk — frameworks like Hermes",
    "auto-write skills with zero human review, and their own threat models call memory",
    "“the largest unbounded attack surface.”",
    "",
    "Here, observation is automatic but behavior change never is.",
], body_size=12)
card(s, Inches(6.85), Inches(4.45), Inches(5.9), Inches(2.1), "Seeded today", [
    "10 corporate + 2 HSA instincts (Ansible standards, GitLab CI security,",
    "PCI controls, secrets patterns, supply-chain hardening) — each with",
    "citation, approver, and confidence on record.",
], accent=GREEN, body_size=12)
footer(s)

# ================================================================ 10 · GENOME
s = new_slide()
header(s, "Architecture strategy", "The genetic architecture — 14 genes, framework-agnostic")
txt(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7), [
    "We extracted the system into a canonical genome (genome.yaml): 14 invariants with conformance tests.",
    "Any framework “expression” of the system is validated against it — the Claude Code plugin is simply the first phenotype.",
], size=13, color=MUTED, space_after=3)
genes = ["G1 lean orchestrator", "G2 specialist roster", "G3 delegation envelope", "G4 verdict contracts",
         "G5 deterministic gate", "G6 bounded remediation", "G7 model tiering", "G8 tool-boundary enforcement",
         "G9 governance ledger", "G10 zone separation", "G11 lazy knowledge", "G12 governed learning",
         "G13 depth-1 delegation", "G14 deterministic artifacts"]
for i, g in enumerate(genes):
    x = Inches(0.6 + (i % 4) * 3.1)
    y = Inches(2.35 + (i // 4) * 0.62)
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.95), Inches(0.52), fill=PANEL)
    try: sp.adjustments[0] = 0.2
    except Exception: pass
    set_shape_text(sp, g, size=10.5, bold=True, color=TEXT)
bars = [("Class A — ports verbatim", "rules · skills content · instinct ledger · merge gate · state store · templates · validators", GREEN, 8.0),
        ("Class B — re-express per framework", "orchestration topology · roster wiring · model tiering · zones", AMBER, 5.4),
        ("Class C — harness-coupled (thin)", "hook bindings · frontmatter · plugin packaging", RED, 2.6)]
y = Inches(5.05)
for t, d, a, wlen in bars:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(wlen), Inches(0.42), fill=a)
    txt(s, Inches(0.72), y + Inches(0.04), Inches(wlen), Inches(0.35), t, size=11, bold=True, color=DARKTXT)
    txt(s, Inches(8.75), y + Inches(0.04), Inches(4.0), Inches(0.4), d, size=9.5, color=MUTED)
    y += Inches(0.56)
txt(s, Inches(0.6), Inches(6.78), Inches(12.2), Inches(0.35),
    "≈ 70% of the system by value is framework-independent today — we are not locked in.",
    size=13, color=TEAL, bold=True)
footer(s)

# ================================================================ 11 · TOP 5 HARNESSES
s = new_slide()
header(s, "Harness strategy", "Top 5 recommended harnesses (ranked)")
recs = [
    ("1", "Claude Code / Claude Agent SDK", "RUN NOW", GREEN,
     "The reference expression — everything on the previous slides is live on it. Agent SDK re-hosts it headlessly (CI, services, GitLab webhooks). Genome fit: 9/14 native, 0 misfits."),
    ("2", "LangGraph", "BUILD NEXT (HSA)", TEAL,
     "Turns our prose orchestration contract into enforced graph structure (gate = code node, remediation = counted cycle) and runs HSA agents on local models natively — the air-gap path our current harness can never provide. Fit: 12/14 native."),
    ("3", "Microsoft Agent Framework 1.0", "ENTERPRISE ALTERNATE", TEAL,
     "AutoGen + Semantic Kernel successor, GA Apr 2026, LTS. Equivalent graph fidelity; OpenTelemetry-native audit feeds the SIEM with the least glue; first-class .NET for the Windows estate. Fit: 11/14 native."),
    ("4", "Anthropic Managed Agents", "WATCH (BETA)", AMBER,
     "Hosted coordinator that enforces our delegation rule at the API level; rubric-graded outcome loops mirror our review gate; vaults + versioned agents are strong. Corporate lane only — inference is never air-gappable. Re-evaluate at GA."),
    ("5", "NanoClaw", "EVALUATE", AMBER,
     "Hardened minimal host (~700 LOC) over the Claude Agent SDK with real container-per-channel isolation — a credible chat-ops front door / constrained host. The only Claw-family member with a defensible security profile."),
]
y = Inches(1.6)
for n, name, tag, a, d in recs:
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.15), Inches(0.98), fill=PANEL)
    try: sp.adjustments[0] = 0.1
    except Exception: pass
    nb = shape(s, MSO_SHAPE.OVAL, Inches(0.75), y + Inches(0.22), Inches(0.55), Inches(0.55), fill=a)
    set_shape_text(nb, n, size=18, bold=True, color=DARKTXT)
    txt(s, Inches(1.5), y + Inches(0.06), Inches(4.6), Inches(0.45), name, size=14.5, bold=True)
    tg = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y + Inches(0.52), Inches(1.95), Inches(0.36), fill=a)
    set_shape_text(tg, tag, size=9.5, bold=True, color=DARKTXT)
    txt(s, Inches(4.7), y + Inches(0.05), Inches(7.9), Inches(0.9), d, size=10.5, color=TEXT,
        anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.06)
footer(s)

# ================================================================ 12 · WHO DIDN'T MAKE IT
s = new_slide()
header(s, "Harness strategy", "Evaluated and excluded — and why")
outs = [
    ("OpenAI Agents SDK", "Its flagship peer-handoff idiom is the opposite of our hub-and-spoke control rule; model tiering and the local lane go through shims. Viable only if the org commits to OpenAI."),
    ("CrewAI", "Fastest roster prototype, but tool-boundary enforcement is advisory callbacks — not a fail-closed gate every call must cross. Disqualifying for a CDE; fine for throwaway prototypes."),
    ("Google ADK", "Solid callbacks and unique A2A interop, but its Gemini/Vertex center of gravity fights a self-hosted, air-gapped estate at every step. Revisit if GCP enters the picture."),
    ("OpenClaw (core)", "≈12% of marketplace skills found malicious (341/2,857); CVSS 8.8 one-click RCE (CVE-2026-25253); its channel breadth is exactly the attack surface a CDE must remove."),
    ("Hermes Agent", "The anti-pattern of our learning gene: ungoverned self-improvement, auto-written skills, memory as an unbounded attack surface. We steal its trace-mining idea for instinct candidates — gated by humans."),
]
y = Inches(1.65)
for name, d in outs:
    card(s, Inches(0.6), y, Inches(12.15), Inches(0.92), name, d, accent=RED, title_size=12.5, body_size=10.5)
    y += Inches(1.0)
txt(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.35),
    "Ideas adopted from the excluded: SKILL.md cross-vendor portability · GEPA-style trace mining (human-gated) · TEE attestation & kernel sandboxing as audit-depth upgrades.",
    size=11, color=MUTED)
footer(s)

# ================================================================ 13 · ROADMAP
s = new_slide()
header(s, "Roadmap", "From proof of concept to production")
cols = [
    ("NOW — unblock the PoC (P0)", GREEN, [
        "• Stand up local model on the PoC box (Ollama + Qwen-coder) — the local lane is wired, endpoint isn't",
        "• Create least-privilege GitLab service accounts (read + branch/MR write only)",
        "• Run /infra-discover against the real estate → knowledge/environment.md (ground truth)",
    ]),
    ("NEXT — prove portability & HSA path", AMBER, [
        "• Prototype the HSA expression in LangGraph: existing perso-* prompts + merge-gate node + local models",
        "• Wire the SIEM endpoint for governance events (PCI CP §6.4 retention)",
        "• Validate against the genome conformance checklist (red-team the DLP / zone gates)",
    ]),
    ("LATER — production hardening", TEAL, [
        "• CPSA design review → HSA deployment unblock (compliance gate, not technical)",
        "• Runner topology split: CI / Deploy / Windows / HSA runners (closes the single-box PCI gap)",
        "• Molecule test retrofit for existing playbooks; periodic environment re-discovery",
    ]),
]
for i, (t, a, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.13)
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.65), Inches(3.9), Inches(4.6), fill=PANEL)
    try: sp.adjustments[0] = 0.04
    except Exception: pass
    shape(s, MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(3.9), Inches(0.55), fill=a)
    txt(s, x + Inches(0.15), Inches(1.7), Inches(3.6), Inches(0.45), t, size=12.5, bold=True, color=DARKTXT)
    txt(s, x + Inches(0.18), Inches(2.4), Inches(3.6), Inches(3.7), items, size=11, space_after=8)
txt(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.5),
    "Decision checkpoint: after P0 lands and CPSA review is scheduled — confirm LangGraph as the HSA expression (alternate: MS Agent Framework).",
    size=12, color=AMBER, bold=True)
footer(s)

# ================================================================ 14 · THE ASK
s = new_slide()
header(s, "Closing", "What we're asking for")
asks = [
    ("Green-light PoC continuation", "Approve the P0 work: local-model hardware on the PoC box, GitLab service accounts, estate discovery run."),
    ("Sponsor the HSA prototype", "2–3 weeks to express the air-gapped lane in LangGraph using assets we already have — development is not CPSA-gated, only deployment is."),
    ("Schedule the CPSA conversation", "HSA deployment is blocked on compliance review, not engineering. Earlier scheduling = earlier value in the zone that matters most."),
]
y = Inches(1.7)
for t, d in asks:
    card(s, Inches(0.6), y, Inches(12.15), Inches(1.15), t, d, accent=TEAL, title_size=15, body_size=12.5)
    y += Inches(1.3)
sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.7), Inches(12.15), Inches(1.1), fill=PANEL2)
try: sp.adjustments[0] = 0.08
except Exception: pass
set_shape_text(sp, [
    ("The one-line summary", {"size": 12, "bold": True, "color": AMBER}),
    ("An agent system a QSA can audit: every action ledgered, every merge computed, every learned behavior human-approved — and 70% portable across the 2026 harness landscape.",
     {"size": 14, "bold": True, "color": TEXT}),
], align=PP_ALIGN.LEFT)
footer(s)

# ================================================================ save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "infra-ops-poc.pptx")
prs.save(out)
print(f"wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
